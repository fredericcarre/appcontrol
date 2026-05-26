#!/usr/bin/env python3
"""
Build AppControl_DORA_Invivoo_Dark.pptx — Option D:
26 slides alternating
  - cream Invivoo illustration (from V1, full-bleed)
  - dense dark-navy content slide in Invivoo brand colors
Keeps the V1 punch and density; only the dense slides are re-skinned
from black/cyan/magenta to Invivoo navy + cream + #4E57AA accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Invivoo DARK palette ─────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x2C, 0x5F)   # primary slide background
NAVY_DEEP  = RGBColor(0x08, 0x1F, 0x42)   # darker accent (code blocks)
NAVY_PANEL = RGBColor(0x1A, 0x38, 0x70)   # card surface
NAVY_LINE  = RGBColor(0x3A, 0x52, 0x94)   # card borders
CREAM      = RGBColor(0xF5, 0xEF, 0xE4)   # primary text
CREAM_MUTE = RGBColor(0xC7, 0xCF, 0xE5)   # secondary text
INVIVOO_BLUE   = RGBColor(0x4E, 0x57, 0xAA)   # primary accent
INVIVOO_BLUE_2 = RGBColor(0x7A, 0x86, 0xD8)   # lighter accent
AMBER      = RGBColor(0xE8, 0xA3, 0x3D)   # big numbers / sanctions
TEAL       = RGBColor(0x5B, 0xC9, 0xA8)   # ok / positive
CRIMSON    = RGBColor(0xE2, 0x5C, 0x5C)   # danger
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PURE_BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

ILLUSTRATIONS = "/tmp/v1-illustrations"

# ── Helpers ──────────────────────────────────────────────────────────

def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False

def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=CREAM, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri", line_spacing=1.20):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
    return tb

def add_rule(slide, left, top, width, *, color=INVIVOO_BLUE, weight=2.0):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_panel(slide, left, top, width, height, *,
              fill=NAVY_PANEL, border=NAVY_LINE, border_weight=0.75,
              radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    if radius:
        shp.adjustments[0] = 0.03
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border_weight > 0:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_weight)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_topbar(slide, kicker):
    add_rule(slide, Inches(0.55), Inches(0.55), Inches(0.55),
             color=INVIVOO_BLUE, weight=2.0)
    add_text(slide, Inches(1.20), Inches(0.40), Inches(8), Inches(0.40),
             kicker, size=10, bold=True, color=INVIVOO_BLUE_2)
    add_text(slide, Inches(11.0), Inches(0.40), Inches(1.8), Inches(0.40),
             "INVIVOO", size=10, bold=True, color=CREAM,
             align=PP_ALIGN.RIGHT)

def add_footer(slide, n, total=26):
    add_rule(slide, Inches(0.55), Inches(7.10), Inches(12.23),
             color=NAVY_LINE, weight=0.5)
    add_text(slide, Inches(0.55), Inches(7.18), Inches(8), Inches(0.30),
             "Invivoo  ·  AppControl  ·  Pitch DORA  ·  Mai 2026  ·  Confidentiel",
             size=8.5, color=CREAM_MUTE)
    add_text(slide, Inches(11.0), Inches(7.18), Inches(1.8), Inches(0.30),
             f"{n} / {total}", size=8.5, color=CREAM_MUTE,
             align=PP_ALIGN.RIGHT)

def add_title(slide, top, title, *, size=30, color=CREAM):
    add_text(slide, Inches(0.55), Inches(top), Inches(12.2), Inches(0.9),
             title, size=size, bold=True, color=color, line_spacing=1.10)

def add_badge(slide, left, top, width, height, label, *, fill=INVIVOO_BLUE,
              text_color=CREAM, size=10.5):
    add_panel(slide, left, top, width, height,
              fill=fill, border=fill, border_weight=0)
    add_text(slide, left, top, width, height, label,
             size=size, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ── Slide builders ───────────────────────────────────────────────────

def illustration_slide(image_path):
    s = prs.slides.add_slide(BLANK)
    # Full-bleed: the V1 illustrations are ~16:9 cream/blue compositions.
    s.shapes.add_picture(image_path, 0, 0, SLIDE_W, SLIDE_H)
    return s

# ────────── 2. PROBLÈME ──────────
def slide_probleme(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "LE PROBLÈME  ·  INVENTAIRE DE LA VÉRITÉ")
    add_title(s, 1.05,
              "L'illusion de la connaissance :", size=32, color=CREAM)
    add_title(s, 1.55,
              "une architecture éclatée et purement déclarative",
              size=24, color=CREAM_MUTE)
    sources = [
        ("CMDB",              "Souvent obsolète, limitée aux briques techniques",  "OBSOLÈTE", CRIMSON),
        ("Runbooks",          "Procédures rédigées hors-sol, rarement à jour",     "OBSOLÈTE", CRIMSON),
        ("Sachants",          "Connaissance tribale dispersée",                     "FRAGILE",  AMBER),
        ("XL Release / XL Deploy", "Vue déploiement, mais pas opérationnelle",     "PARTIEL",  AMBER),
        ("Schémas Visio",     "Intentions figées, déconnectées de la réalité",     "FIGÉ",     AMBER),
        ("Référentiel de flux", "Vérité terrain brute, jamais capitalisée",        "BRUT",     INVIVOO_BLUE_2),
    ]
    cols = 3; cw = Inches(4.05); ch = Inches(1.55)
    left0 = Inches(0.55); top0 = Inches(2.70)
    gx = Inches(0.10); gy = Inches(0.20)
    for i, (name, desc, tag, tagcol) in enumerate(sources):
        r, c = divmod(i, cols)
        x = left0 + (cw + gx) * c
        y = top0  + (ch + gy) * r
        add_panel(s, x, y, cw, ch)
        add_text(s, x + Inches(0.20), y + Inches(0.15),
                 cw - Inches(0.40), Inches(0.40),
                 name, size=14, bold=True, color=CREAM)
        add_text(s, x + Inches(0.20), y + Inches(0.60),
                 cw - Inches(0.40), Inches(0.70),
                 desc, size=11, color=CREAM_MUTE, line_spacing=1.25)
        add_badge(s, x + cw - Inches(1.10), y + Inches(1.10),
                  Inches(0.95), Inches(0.32), tag,
                  fill=tagcol, text_color=NAVY, size=9)
    # Alert
    add_panel(s, Inches(0.55), Inches(6.10),
              Inches(12.23), Inches(0.85),
              fill=NAVY_DEEP, border=CRIMSON, border_weight=2.0)
    add_text(s, Inches(0.85), Inches(6.18), Inches(0.30), Inches(0.70),
             "!", size=24, bold=True, color=CRIMSON, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.25), Inches(6.18), Inches(11.4), Inches(0.70),
             ["Aucune source ne décrit comment l'application tourne réellement en production.",
              "Sans cette vérité consolidée, la reconstruction est impossible."],
             size=12.5, color=CREAM, line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 4. CHOC DORA ──────────
def slide_dora(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "RÉGLEMENTATION  ·  DORA")
    add_title(s, 1.05,
              "Le choc DORA :", size=32, color=CREAM)
    add_title(s, 1.55,
              "l'obligation légale de prouver la reconstruction",
              size=24, color=CREAM_MUTE)
    # Big sanctions
    panels = [
        ("SANCTION ENTREPRISE", "2 %",  "du chiffre d'affaires annuel mondial"),
        ("SANCTION DIRIGEANT",  "1 M€", "d'amende personnelle"),
    ]
    pw = Inches(3.20); ph = Inches(2.40)
    pl = Inches(0.55); pg = Inches(0.20)
    for i, (label, big, sub) in enumerate(panels):
        x = pl + (pw + pg) * i
        y = Inches(2.70)
        add_panel(s, x, y, pw, ph,
                  fill=NAVY_PANEL, border=AMBER, border_weight=2.0)
        add_text(s, x + Inches(0.25), y + Inches(0.20),
                 pw - Inches(0.5), Inches(0.4),
                 label, size=10, bold=True, color=AMBER)
        add_text(s, x + Inches(0.25), y + Inches(0.65),
                 pw - Inches(0.5), Inches(1.30),
                 big, size=72, bold=True, color=AMBER, line_spacing=1.0)
        add_text(s, x + Inches(0.25), y + Inches(1.95),
                 pw - Inches(0.5), Inches(0.4),
                 sub, size=12, color=CREAM)
    # Right side: articles
    arts = [
        ("Art. 11 & 12", "Procédures de reconstruction après corruption majeure, testées annuellement."),
        ("Art. 25",      "Tests de scénarios cyber chronométrés."),
        ("Art. 16",      "Registre inaltérable des actions de récupération."),
    ]
    ax = pl + (pw + pg) * 2 + Inches(0.20)
    aw = SLIDE_W - ax - Inches(0.55)
    ay = Inches(2.70)
    for i, (art, body) in enumerate(arts):
        y = ay + Inches(0.80) * i
        add_badge(s, ax, y, Inches(1.20), Inches(0.32), art,
                  fill=INVIVOO_BLUE, size=10)
        add_text(s, ax + Inches(1.35), y - Inches(0.02),
                 aw - Inches(1.35), Inches(0.75),
                 body, size=12, color=CREAM, line_spacing=1.25)
    # Bottom keyline
    add_panel(s, Inches(0.55), Inches(5.50),
              Inches(12.23), Inches(1.40),
              fill=AMBER, border=AMBER, border_weight=0)
    add_text(s, Inches(0.55), Inches(5.50),
             Inches(12.23), Inches(1.40),
             ["DORA exige de prouver la capacité de reprise,",
              "pas seulement de la documenter."],
             size=22, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.25)
    add_footer(s, n)

# ────────── 6. PIÈGE STRATÉGIQUE ──────────
def slide_piege(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "ANALYSE COMPARATIVE  ·  POURQUOI APPCONTROL")
    add_title(s, 1.05,
              "Le piège stratégique :", size=32, color=CREAM)
    add_title(s, 1.55,
              "pourquoi une CMDB parfaite ne suffit pas",
              size=24, color=CREAM_MUTE)
    # Comparison table
    rows = [
        ("Time-to-value",       "3 à 5 ans (fort risque d'échec)",       "Semaines par application"),
        ("Nature de la donnée", "Déclarative (saisie humaine)",          "Observée (réconciliation temps réel)"),
        ("Problème résolu",     "Qualité par silo",                       "Exploitation inter-silos"),
        ("Capacité de Rebuild", "NON · un référentiel ne reconstruit pas","OUI · c'est sa fonction native"),
    ]
    col_w = [Inches(3.20), Inches(4.50), Inches(4.50)]
    left = Inches(0.55); top = Inches(2.70); row_h = Inches(0.70)
    headers = ["CRITÈRE", "Refonte classique des référentiels",
               "AppControl  ·  Outil transverse"]
    head_cols = [None, CRIMSON, TEAL]
    for j, (h, w) in enumerate(zip(headers, col_w)):
        x = left + sum(col_w[:j], Emu(0))
        if j == 0:
            add_text(s, x + Inches(0.20), top + Inches(0.15),
                     w - Inches(0.4), row_h - Inches(0.3),
                     h, size=11, bold=True, color=INVIVOO_BLUE_2)
        else:
            add_panel(s, x, top, w, row_h,
                      fill=head_cols[j], border=head_cols[j], border_weight=0)
            add_text(s, x, top, w, row_h,
                     h, size=12, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, row in enumerate(rows):
        ry = top + row_h * (i + 1)
        if i % 2 == 0:
            add_panel(s, left, ry, sum(col_w, Emu(0)), row_h,
                      fill=NAVY_PANEL, border=NAVY_PANEL, border_weight=0)
        for j, (cell, w) in enumerate(zip(row, col_w)):
            x = left + sum(col_w[:j], Emu(0))
            if j == 0:
                add_text(s, x + Inches(0.20), ry, w - Inches(0.40), row_h,
                         cell, size=12, bold=True, color=CREAM,
                         anchor=MSO_ANCHOR.MIDDLE)
            else:
                col = CRIMSON if j == 1 else TEAL
                mark = "✕" if j == 1 else "✓"
                add_text(s, x + Inches(0.20), ry, Inches(0.40), row_h,
                         mark, size=16, bold=True, color=col,
                         anchor=MSO_ANCHOR.MIDDLE)
                add_text(s, x + Inches(0.65), ry, w - Inches(0.85), row_h,
                         cell, size=12, color=CREAM, anchor=MSO_ANCHOR.MIDDLE)
    # Footer alert
    add_panel(s, Inches(0.55), Inches(6.10),
              Inches(12.23), Inches(0.85),
              fill=NAVY_DEEP, border=AMBER, border_weight=2.0)
    add_text(s, Inches(0.55), Inches(6.10),
             Inches(12.23), Inches(0.85),
             "Une CMDB parfaite ne peut pas redémarrer une application. DORA exige l'exécution, pas l'inventaire.",
             size=13.5, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 8. PROPOSITION DE VALEUR ──────────
def slide_valeur(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "PROPOSITION DE VALEUR  ·  APPCONTROL")
    add_title(s, 1.05,
              "De la donnée morte au plan opérable", size=34)
    add_title(s, 1.60,
              "Trois temps  ·  un seul fil conducteur",
              size=15, color=CREAM_MUTE)
    steps = [
        ("01", "Capter",
         "Récupérer la donnée existante (CMDB, XLR, XLD, flux, logs).",
         "Sans imposer de nouvelle saisie.",
         "ENTONNOIR DE COLLECTE",
         "CMDB · XLR · XLD · Flux · Logs"),
        ("02", "Réconcilier",
         "Croiser les sources, lever les contradictions.",
         "Produire une carte vivante fidèle à la production.",
         "MOTEUR DE RÉCONCILIATION",
         "AppControl"),
        ("03", "Exploiter",
         "Utiliser cette carte comme moteur d'exécution.",
         "Démarrage séquencé, rebuild, bascule DR, audit.",
         "PLAN D'EXÉCUTION",
         "Séquence · Rebuild · DR · Audit"),
    ]
    sw = Inches(4.00); sh = Inches(3.90)
    sl = Inches(0.55); sg = Inches(0.30)
    for i, (n_, h, d1, d2, kicker, footer) in enumerate(steps):
        x = sl + (sw + sg) * i
        y = Inches(2.50)
        add_panel(s, x, y, sw, sh)
        # Number
        add_text(s, x + Inches(0.30), y + Inches(0.20),
                 Inches(1.5), Inches(0.8),
                 n_, size=44, bold=True, color=INVIVOO_BLUE_2,
                 line_spacing=1.0)
        add_text(s, x + Inches(0.30), y + Inches(1.10),
                 sw - Inches(0.6), Inches(0.7),
                 h, size=26, bold=True, color=CREAM)
        add_rule(s, x + Inches(0.30), y + Inches(1.95),
                 Inches(1.20), color=INVIVOO_BLUE, weight=2.0)
        add_text(s, x + Inches(0.30), y + Inches(2.10),
                 sw - Inches(0.6), Inches(0.75),
                 d1, size=12.5, color=CREAM, line_spacing=1.35)
        add_text(s, x + Inches(0.30), y + Inches(2.80),
                 sw - Inches(0.6), Inches(0.55),
                 d2, size=11, color=CREAM_MUTE, line_spacing=1.3)
        add_rule(s, x + Inches(0.30), y + Inches(3.30),
                 sw - Inches(0.6), color=NAVY_LINE, weight=0.75)
        add_text(s, x + Inches(0.30), y + Inches(3.35),
                 sw - Inches(0.6), Inches(0.30),
                 kicker, size=9, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x + Inches(0.30), y + Inches(3.60),
                 sw - Inches(0.6), Inches(0.30),
                 footer, size=10.5, color=CREAM)
    # Bottom strip
    add_panel(s, Inches(0.55), Inches(6.60),
              Inches(12.23), Inches(0.40),
              fill=INVIVOO_BLUE, border=INVIVOO_BLUE, border_weight=0)
    add_text(s, Inches(0.55), Inches(6.60),
             Inches(12.23), Inches(0.40),
             "Pas de redocumentation. Pas de big bang. La carte est observée, validée et exécutée — au rythme de chaque application.",
             size=11, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 10. PHASE 1 — X-Ray ──────────
def slide_phase1(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "PHASE 1  ·  CAPTATION  ·  AS-RUN")
    add_title(s, 1.05,
              "Révéler l'architecture « As-Run »", size=30)
    add_title(s, 1.55,
              "sans redocumenter", size=30, color=CREAM_MUTE)
    # 3 sources
    sources = [
        ("DONNÉES FROIDES", "CMDB",                 "Serveurs, middlewares"),
        ("DONNÉES CHAUDES", "XL Deploy / XL Release","Manifestes, mapping composant,\nordre de déploiement"),
        ("VÉRITÉ TERRAIN",  "Référentiels de flux", "Flux applicatifs, supervision"),
    ]
    cw = Inches(2.95); ch = Inches(1.85)
    cl = Inches(0.55); cg = Inches(0.20)
    cy = Inches(2.55)
    for i, (k, h, d) in enumerate(sources):
        x = cl + (cw + cg) * i
        add_panel(s, x, cy, cw, ch)
        add_text(s, x + Inches(0.20), cy + Inches(0.18),
                 cw - Inches(0.4), Inches(0.30),
                 k, size=9, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x + Inches(0.20), cy + Inches(0.50),
                 cw - Inches(0.4), Inches(0.50),
                 h, size=15, bold=True, color=CREAM)
        add_text(s, x + Inches(0.20), cy + Inches(1.05),
                 cw - Inches(0.4), Inches(0.75),
                 d, size=10.5, color=CREAM_MUTE, line_spacing=1.3)
    # JSON outputs (right side)
    json_x = cl + (cw + cg) * 3
    json_w = SLIDE_W - json_x - Inches(0.55)
    json_blocks = [
        ("Topologie d'application",
         '{"app":"core","status":"active",\n "dependencies":["db_cluster_1","auth_service"]}'),
        ("Composants & versions",
         '{"component":"payment_gateway",\n "version":"2.3.1","type":"microservice"}'),
        ("Flux & infrastructure",
         '{"flow":"user_login","source":"web_ui",\n "target":"auth_api","host":"server_x"}'),
    ]
    bh = Inches(0.62)
    by = cy
    for i, (label, code) in enumerate(json_blocks):
        y = by + Inches(0.65) * i
        add_text(s, json_x, y, json_w, Inches(0.30),
                 label, size=10, bold=True, color=INVIVOO_BLUE_2)
        add_panel(s, json_x, y + Inches(0.28), json_w, Inches(0.55),
                  fill=NAVY_DEEP, border=NAVY_LINE, border_weight=0.5)
        add_text(s, json_x + Inches(0.10), y + Inches(0.32),
                 json_w - Inches(0.20), Inches(0.50),
                 code, size=9, color=TEAL, font="Consolas",
                 line_spacing=1.20)
    # X-Ray scanner bar
    add_panel(s, Inches(2.55), Inches(4.85),
              Inches(8.20), Inches(0.85),
              fill=INVIVOO_BLUE, border=INVIVOO_BLUE, border_weight=0)
    add_text(s, Inches(2.55), Inches(4.85), Inches(8.20), Inches(0.85),
             "APPCONTROL X-RAY SCANNER   ·   Capter · Réconcilier · Générer",
             size=14, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom keyline
    add_panel(s, Inches(0.55), Inches(6.00),
              Inches(12.23), Inches(0.95),
              fill=NAVY_DEEP, border=AMBER, border_weight=1.5)
    add_text(s, Inches(0.75), Inches(6.00), Inches(11.93), Inches(0.95),
             ["CLÉ DE VOÛTE — la donnée XL Deploy / Release est fraîche (regénérée à chaque release) et opérationnelle.",
              "AppControl produit un premier JSON de l'application tel qu'elle tourne aujourd'hui, soumis à la validation des sachants."],
             size=11, color=CREAM, line_spacing=1.30,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 12. PHASE 2 — Carte télécommande ──────────
def slide_phase2(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "PHASE 2  ·  EXPLOITATION")
    add_title(s, 1.05,
              "La carte devient la télécommande", size=30)
    add_title(s, 1.55,
              "de l'application", size=30, color=CREAM_MUTE)
    # Top-left: 8 component squares (Visual Map Controller)
    map_x = Inches(0.55); map_y = Inches(2.55)
    map_w = Inches(5.80); map_h = Inches(2.80)
    add_panel(s, map_x, map_y, map_w, map_h)
    add_text(s, map_x + Inches(0.20), map_y + Inches(0.15),
             map_w - Inches(0.4), Inches(0.30),
             "VISUAL MAP CONTROLLER",
             size=9.5, bold=True, color=INVIVOO_BLUE_2)
    comps = ["DB", "SERVICE", "API", "NODE",
             "CLUSTER", "GATEWAY", "CACHE", "WORKER"]
    inner_x = map_x + Inches(0.25); inner_y = map_y + Inches(0.60)
    bw = Inches(1.30); bh = Inches(0.85)
    bgx = Inches(0.07); bgy = Inches(0.20)
    for i, label in enumerate(comps):
        r, c = divmod(i, 4)
        x = inner_x + (bw + bgx) * c
        y = inner_y + (bh + bgy) * r
        add_panel(s, x, y, bw, bh,
                  fill=NAVY_DEEP, border=INVIVOO_BLUE, border_weight=1.0)
        add_text(s, x, y + Inches(0.10), bw, Inches(0.30),
                 label, size=11, bold=True, color=CREAM,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.42), bw, Inches(0.40),
                 "START · STOP · CHECK", size=7, color=TEAL,
                 align=PP_ALIGN.CENTER)
    # restart-branch indicator
    add_panel(s, map_x + Inches(0.25), map_y + Inches(2.40),
              Inches(5.30), Inches(0.30),
              fill=AMBER, border=AMBER, border_weight=0)
    add_text(s, map_x + Inches(0.25), map_y + Inches(2.40),
             Inches(5.30), Inches(0.30),
             "↻  RESTART BRANCH",
             size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Top-right: GitOps PR
    pr_x = map_x + map_w + Inches(0.20)
    pr_w = SLIDE_W - pr_x - Inches(0.55)
    pr_y = map_y
    add_panel(s, pr_x, pr_y, pr_w, map_h, fill=NAVY_DEEP)
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(0.15),
             pr_w - Inches(0.4), Inches(0.30),
             "GITOPS CODE INTERFACE  ·  PULL REQUEST",
             size=9.5, bold=True, color=INVIVOO_BLUE_2)
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(0.50),
             pr_w - Inches(0.4), Inches(0.30),
             "PR #452  ·  Targeted Restart Automation",
             size=10, bold=True, color=CREAM)
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(0.85),
             pr_w - Inches(0.4), Inches(0.30),
             "Commit: Fix 'pink branch' dependency chain",
             size=9, color=CREAM_MUTE, font="Consolas")
    # before/after snippets
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(1.20),
             Inches(0.6), Inches(0.30),
             "− BEFORE", size=9, bold=True, color=CRIMSON, font="Consolas")
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(1.45),
             pr_w - Inches(0.4), Inches(0.85),
             '"status": "error_state",\n"active_flow": false,\n"action":       "—"',
             size=9, color=CRIMSON, font="Consolas", line_spacing=1.25)
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(2.20),
             Inches(0.6), Inches(0.30),
             "+ AFTER",  size=9, bold=True, color=TEAL, font="Consolas")
    add_text(s, pr_x + Inches(0.20), pr_y + Inches(2.45),
             pr_w - Inches(0.4), Inches(0.4),
             '"status": "ok", "active_flow": true,\n"action": "restart_targeted"',
             size=9, color=TEAL, font="Consolas", line_spacing=1.25)
    # Bottom: 4 capabilities
    caps = [
        ("01", "Démarrage et arrêt séquencés",
         "respectant le graphe de dépendances (DAG)"),
        ("02", "Restart ciblé sur branche en erreur",
         "(« pink branch »)"),
        ("03", "Intégration native avec les schedulers",
         "Control-M, AutoSys, $U, TWS"),
        ("04", "Bascule DR orchestrée",
         "en 6 phases avec rollback"),
    ]
    cw = Inches(2.95); ch = Inches(1.30)
    cl = Inches(0.55); cg = Inches(0.20)
    cy = Inches(5.55)
    for i, (n_, h, d) in enumerate(caps):
        x = cl + (cw + cg) * i
        add_panel(s, x, cy, cw, ch)
        add_text(s, x + Inches(0.15), cy + Inches(0.10),
                 Inches(0.5), Inches(0.45),
                 n_, size=18, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x + Inches(0.65), cy + Inches(0.13),
                 cw - Inches(0.85), Inches(0.6),
                 h, size=11, bold=True, color=CREAM, line_spacing=1.20)
        add_text(s, x + Inches(0.15), cy + Inches(0.78),
                 cw - Inches(0.30), Inches(0.45),
                 d, size=9.5, color=CREAM_MUTE, line_spacing=1.30)
    add_footer(s, n)

# ────────── 14. PHASE 3 — Preuve du rebuild ──────────
def slide_phase3(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "PHASE 3  ·  PREUVE DU REBUILD")
    add_title(s, 1.05,
              "Prouver le Rebuild —", size=30)
    add_title(s, 1.55,
              "la conformité DORA transformée en code",
              size=24, color=CREAM_MUTE)
    # 3 columns
    # Left: input map_config.json
    lx = Inches(0.55); ly = Inches(2.65); lw = Inches(3.40); lh = Inches(3.20)
    add_panel(s, lx, ly, lw, lh, fill=NAVY_DEEP)
    add_text(s, lx + Inches(0.20), ly + Inches(0.18),
             lw - Inches(0.4), Inches(0.30),
             "INPUT  ·  CARTE EXÉCUTABLE",
             size=9.5, bold=True, color=INVIVOO_BLUE_2)
    add_text(s, lx + Inches(0.20), ly + Inches(0.55),
             lw - Inches(0.4), Inches(2.20),
             '"application": "core-banking",\n"dag": {\n  "db-primary":  ["start"],\n  "cache-layer": ["start"],\n  "api-gateway": ["start"],\n  "scheduler":   ["start"]\n},\n"rebuild_steps": [...],\n"validation":    "auto"',
             size=9, color=TEAL, font="Consolas", line_spacing=1.40)
    add_text(s, lx + Inches(0.20), ly + lh - Inches(0.40),
             lw - Inches(0.4), Inches(0.30),
             "map_config.json  ·  GitOps",
             size=9, color=CREAM_MUTE, font="Consolas")
    # Middle: 2 voies
    mx = lx + lw + Inches(0.20)
    mw = Inches(5.20)
    my = ly
    add_text(s, mx, my - Inches(0.05),
             mw, Inches(0.30),
             "VALIDATION  ·  2 VOIES",
             size=9.5, bold=True, color=INVIVOO_BLUE_2)
    voie_h = Inches(1.50); voie_gap = Inches(0.15)
    for i, (badge, h, body, bene) in enumerate([
        ("VOIE A", "Le Dry-Run (Simulation)",
         "Appel API simulant l'ordre du DAG et la résolution des commandes — sans exécution réelle.",
         "BÉNÉFICE  ·  Risque zéro, validation continue de la logique."),
        ("VOIE B", "Le Drill Réel (Non-Prod)",
         "Le moteur exécute la reconstruction réelle sur un site Staging ou DR.",
         "BÉNÉFICE  ·  Preuve d'exécution chronométrée, livrable d'audit DORA."),
    ]):
        y = my + Inches(0.30) + (voie_h + voie_gap) * i
        add_panel(s, mx, y, mw, voie_h)
        add_badge(s, mx + Inches(0.20), y + Inches(0.15),
                  Inches(1.00), Inches(0.32), badge, fill=INVIVOO_BLUE, size=10)
        add_text(s, mx + Inches(1.30), y + Inches(0.12),
                 mw - Inches(1.50), Inches(0.40),
                 h, size=13, bold=True, color=CREAM)
        add_text(s, mx + Inches(0.20), y + Inches(0.55),
                 mw - Inches(0.40), Inches(0.55),
                 body, size=10.5, color=CREAM_MUTE, line_spacing=1.30)
        add_text(s, mx + Inches(0.20), y + voie_h - Inches(0.32),
                 mw - Inches(0.40), Inches(0.28),
                 bene, size=9.5, bold=True, color=TEAL)
    # Right: DORA stamp
    rx = mx + mw + Inches(0.20)
    rw = SLIDE_W - rx - Inches(0.55)
    add_panel(s, rx, ly, rw, lh,
              fill=NAVY_PANEL, border=TEAL, border_weight=2.0)
    add_text(s, rx + Inches(0.20), ly + Inches(0.20),
             rw - Inches(0.4), Inches(0.30),
             "OUTPUT  ·  PREUVE DORA",
             size=9.5, bold=True, color=INVIVOO_BLUE_2)
    add_text(s, rx, ly + Inches(0.80),
             rw, Inches(1.20),
             "✓", size=84, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(s, rx + Inches(0.20), ly + Inches(2.10),
             rw - Inches(0.4), Inches(0.40),
             "Conformité DORA",
             size=14, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.20), ly + Inches(2.50),
             rw - Inches(0.4), Inches(0.40),
             "GARANTIE",
             size=10, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER)
    # 4 stamps row
    stamps = ["✓ Audit Log", "✓ RTR mesuré",
              "✓ Art. 11 · 12", "✓ Art. 25"]
    sw = Inches(2.95); sh = Inches(0.50)
    sl = Inches(0.55); sg = Inches(0.10)
    sy = Inches(6.00)
    for i, txt in enumerate(stamps):
        x = sl + (sw + sg) * i
        add_panel(s, x, sy, sw, sh,
                  fill=NAVY_DEEP, border=TEAL, border_weight=1.0)
        add_text(s, x, sy, sw, sh, txt,
                 size=11, bold=True, color=TEAL,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Mécanique d'audit
    add_text(s, Inches(0.55), Inches(6.60),
             Inches(12.23), Inches(0.40),
             "MÉCANIQUE D'AUDIT  ·  Le système chronomètre l'exécution et génère un RTR (Recovery Time for Rebuild) réel. Les résultats alimentent l'Audit Log exigé par DORA.",
             size=10, color=CREAM_MUTE, align=PP_ALIGN.CENTER)
    add_footer(s, n)

# ────────── 16. ESCALIER ──────────
def slide_escalier(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "DÉPLOIEMENT  ·  MONTÉE EN PUISSANCE PROGRESSIVE")
    add_title(s, 1.05,
              "L'escalier de la confiance :", size=30)
    add_title(s, 1.55,
              "une montée en puissance sans risque",
              size=24, color=CREAM_MUTE)
    phases = [
        ("PHASE A", "Étapes 1 → 3",
         "Connexion · Réconciliation · Validation",
         "Inventaire transverse, aucune exécution",
         "LECTURE SEULE  ·  RISQUE NUL", TEAL),
        ("PHASE B", "Étapes 4 → 5",
         "Advisory & Diagnostic",
         "Détection des dérives, observation passive",
         "OBSERVATION  ·  RISQUE MAÎTRISÉ", AMBER),
        ("PHASE C", "Étapes 6 → 9",
         "Opérations · Rebuild · DORA Ready",
         "Audit prouvé · conformité activée",
         "EXÉCUTION CONTRÔLÉE  ·  AUDITABLE", INVIVOO_BLUE_2),
    ]
    pw = Inches(4.00); ph = Inches(3.40)
    pl = Inches(0.55); pg = Inches(0.30)
    for i, (badge, etapes, h, d, footer_strip, footer_col) in enumerate(phases):
        x = pl + (pw + pg) * i
        # staircase
        y = Inches(3.20) - Inches(0.30 * i)
        add_panel(s, x, y, pw, ph)
        add_badge(s, x + Inches(0.30), y + Inches(0.30),
                  Inches(1.30), Inches(0.40),
                  badge, fill=INVIVOO_BLUE, size=11)
        add_text(s, x + Inches(0.30), y + Inches(0.95),
                 pw - Inches(0.6), Inches(0.4),
                 etapes, size=11, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x + Inches(0.30), y + Inches(1.35),
                 pw - Inches(0.6), Inches(0.7),
                 h, size=15, bold=True, color=CREAM, line_spacing=1.20)
        add_text(s, x + Inches(0.30), y + Inches(2.20),
                 pw - Inches(0.6), Inches(0.80),
                 d, size=11.5, color=CREAM_MUTE, line_spacing=1.30)
        add_panel(s, x + Inches(0.30), y + ph - Inches(0.50),
                  pw - Inches(0.60), Inches(0.32),
                  fill=footer_col, border=footer_col, border_weight=0)
        add_text(s, x + Inches(0.30), y + ph - Inches(0.50),
                 pw - Inches(0.6), Inches(0.32),
                 footer_strip, size=9, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Bottom note
    add_text(s, Inches(0.55), Inches(6.85),
             Inches(12.2), Inches(0.30),
             "Le profil de risque reste maîtrisé à chaque étape : l'exécution n'arrive qu'après validation complète de la carte.",
             size=11, color=CREAM_MUTE, align=PP_ALIGN.CENTER)
    add_footer(s, n)

# ────────── 18. GOUVERNANCE ──────────
def slide_gouvernance(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "GOUVERNANCE  ·  SÉCURITÉ PAR CONCEPTION")
    add_title(s, 1.05,
              "Sécurité par conception :", size=30)
    add_title(s, 1.55,
              "les prods gardent un contrôle absolu",
              size=24, color=CREAM_MUTE)
    features = [
        ("1", "RBAC Granulaire",
         "Droits stricts par application — chaque rôle a un périmètre exact.",
         "view < operate < edit < manage < owner"),
        ("2", "Validation par Pull Request",
         "Les actions critiques exigent une PR mergée et revue par les sachants Production.",
         "Aucun script ad hoc."),
        ("3", "Advisory Mode & Dry-Run",
         "Simulation systématique avant toute exécution réelle. Le moteur calcule, propose, mais n'agit pas sans validation.",
         ""),
        ("4", "Audit Log DORA (Append-Only)",
         "Registre inaltérable des actions critiques.",
         "action_log · state_transitions · switchover_log"),
    ]
    cw = Inches(2.95); ch = Inches(2.90)
    cl = Inches(0.55); cg = Inches(0.20)
    for i, (n_, h, d, mono) in enumerate(features):
        x = cl + (cw + cg) * i
        y = Inches(2.55)
        add_panel(s, x, y, cw, ch)
        add_text(s, x + Inches(0.25), y + Inches(0.18),
                 Inches(0.5), Inches(0.5),
                 n_, size=24, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x + Inches(0.25), y + Inches(0.80),
                 cw - Inches(0.5), Inches(0.8),
                 h, size=13, bold=True, color=CREAM, line_spacing=1.25)
        add_text(s, x + Inches(0.25), y + Inches(1.60),
                 cw - Inches(0.5), Inches(0.95),
                 d, size=10.5, color=CREAM_MUTE, line_spacing=1.30)
        if mono:
            add_text(s, x + Inches(0.25), y + Inches(2.55),
                     cw - Inches(0.5), Inches(0.30),
                     mono, size=9, color=TEAL, font="Consolas")
    # Principle band
    add_panel(s, Inches(0.55), Inches(5.55),
              Inches(12.23), Inches(0.80),
              fill=INVIVOO_BLUE, border=INVIVOO_BLUE, border_weight=0)
    add_text(s, Inches(0.55), Inches(5.55),
             Inches(12.23), Inches(0.80),
             ["AppControl ne se substitue jamais aux équipes Production.",
              "Chaque action sensible est tracée, simulée, validée et autorisée — le contrôle reste dans les mains des sachants."],
             size=12, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.30)
    # JAMAIS callout
    add_panel(s, Inches(0.55), Inches(6.50),
              Inches(12.23), Inches(0.50),
              fill=NAVY_DEEP, border=AMBER, border_weight=1.5)
    add_text(s, Inches(0.55), Inches(6.50),
             Inches(12.23), Inches(0.50),
             "!  AUCUN UPDATE  ·  AUCUN DELETE  ·  JAMAIS.",
             size=12, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 20. EFFET RÉSEAU ──────────
def slide_reseau(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "EFFET RÉSEAU  ·  CERCLE VERTUEUX")
    add_title(s, 1.05,
              "L'effet réseau :", size=32)
    add_title(s, 1.55,
              "guérir l'écosystème et enrichir les référentiels",
              size=24, color=CREAM_MUTE)
    cards = [
        ("1", "MUTUALISATION", "Un point d'intégration unique",
         "Un seul point d'intégration pour les schedulers et la supervision. La 10ᵉ application coûte beaucoup moins cher à embarquer que la 1ʳᵉ."),
        ("2", "CERCLE VERTUEUX", "Les référentiels se ré-alimentent",
         "La topologie observée par AppControl alimente la CMDB en retour. On casse définitivement le cycle des référentiels obsolètes."),
    ]
    pw = Inches(5.95); ph = Inches(3.30)
    pl = Inches(0.55); pg = Inches(0.30)
    for i, (num, kicker, h, d) in enumerate(cards):
        x = pl + (pw + pg) * i
        y = Inches(2.50)
        add_panel(s, x, y, pw, ph)
        add_text(s, x + Inches(0.30), y + Inches(0.20),
                 Inches(0.7), Inches(0.7),
                 num, size=40, bold=True, color=INVIVOO_BLUE_2, line_spacing=1.0)
        add_text(s, x + Inches(0.30), y + Inches(1.10),
                 pw - Inches(0.6), Inches(0.35),
                 kicker, size=10, bold=True, color=AMBER)
        add_text(s, x + Inches(0.30), y + Inches(1.50),
                 pw - Inches(0.6), Inches(0.6),
                 h, size=20, bold=True, color=CREAM, line_spacing=1.20)
        add_text(s, x + Inches(0.30), y + Inches(2.20),
                 pw - Inches(0.6), Inches(1.0),
                 d, size=12.5, color=CREAM_MUTE, line_spacing=1.35)
    # Promise
    add_panel(s, Inches(0.55), Inches(6.10),
              Inches(12.23), Inches(0.85),
              fill=INVIVOO_BLUE, border=INVIVOO_BLUE, border_weight=0)
    add_text(s, Inches(0.55), Inches(6.10),
             Inches(12.23), Inches(0.85),
             "PROMESSE DU DÉPLOIEMENT  ·  Chaque nouvelle application déployée renforce le réseau, alimente la CMDB et réduit le coût d'embarquement des suivantes.",
             size=12, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 22. ADOPTION ──────────
def slide_adoption(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "ADOPTION  ·  VALEUR PAR LIGNE DE FRONT")
    add_title(s, 1.05,
              "Une valeur immédiate et spécifique", size=30)
    add_title(s, 1.55,
              "pour chaque ligne de front",
              size=24, color=CREAM_MUTE)
    lines = [
        ("STRATÉGIQUE", "Direction & Gouvernance",
         "GRC · Pilotage · Conformité",
         ["Conformité DORA auditable et prouvable.",
          "Risque personnel et financier maîtrisé (amendes évitées).",
          "DR mesurable (RTR/RTO réels tracés)."]),
        ("TACTIQUE", "Production & SRE",
         "Exploitation · Run · Continuité",
         ["Diagnostic 3 niveaux pour une détection précoce.",
          "Start/Stop automatisé (fin des scripts shell éparpillés).",
          "Bascule DR testable sans drame."]),
        ("OPÉRATIONNEL", "Équipes Applicatives",
         "Développement · Architecture · Delivery",
         ["La carte remplace la documentation morte.",
          "Évolution via Pull Request (IaC-friendly).",
          "Onboarding accéléré des nouveaux développeurs."]),
    ]
    cw = Inches(4.00); ch = Inches(3.65)
    cl = Inches(0.55); cg = Inches(0.30)
    for i, (kicker, h, sub, bullets) in enumerate(lines):
        x = cl + (cw + cg) * i
        y = Inches(2.45)
        add_panel(s, x, y, cw, ch)
        add_badge(s, x + Inches(0.30), y + Inches(0.30),
                  Inches(1.80), Inches(0.40),
                  kicker, fill=INVIVOO_BLUE, size=10.5)
        add_text(s, x + Inches(0.30), y + Inches(0.85),
                 cw - Inches(0.6), Inches(0.6),
                 h, size=18, bold=True, color=CREAM)
        add_text(s, x + Inches(0.30), y + Inches(1.50),
                 cw - Inches(0.6), Inches(0.40),
                 sub, size=10.5, bold=True, color=INVIVOO_BLUE_2)
        for j, b in enumerate(bullets):
            by = y + Inches(2.00) + Inches(0.50) * j
            add_text(s, x + Inches(0.30), by, Inches(0.30), Inches(0.40),
                     "✓", size=14, bold=True, color=TEAL)
            add_text(s, x + Inches(0.65), by,
                     cw - Inches(0.95), Inches(0.50),
                     b, size=11.5, color=CREAM, line_spacing=1.30)
    # Shared-value strip
    add_panel(s, Inches(0.55), Inches(6.30),
              Inches(12.23), Inches(0.65),
              fill=AMBER, border=AMBER, border_weight=0)
    add_text(s, Inches(0.55), Inches(6.30),
             Inches(12.23), Inches(0.65),
             "VALEUR PARTAGÉE  ·  Chaque ligne de front trouve son intérêt dès la première application. AppControl n'est pas un projet IT, c'est un outil quotidien pour 3 publics différents.",
             size=11.5, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ────────── 24. CONCLUSION ──────────
def slide_conclusion(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_topbar(s, "CONCLUSION  ·  NEXT STEP")
    # Big quote
    add_text(s, Inches(0.55), Inches(1.30),
             Inches(12.2), Inches(2.10),
             ["L'exécution est la seule preuve",
              "de la résilience."],
             size=52, bold=True, color=CREAM, line_spacing=1.08)
    # Keyline
    add_panel(s, Inches(0.55), Inches(3.50),
              Inches(12.23), Inches(0.90),
              fill=AMBER, border=AMBER, border_weight=0)
    add_text(s, Inches(0.55), Inches(3.50),
             Inches(12.23), Inches(0.90),
             "Le rebuild n'est pas une option, c'est une obligation DORA.",
             size=20, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 3 pillars
    pillars = [
        ("01", "Immédiateté",
         "Valeur générée dès la première captation, sans attendre des années."),
        ("02", "Sécurité",
         "Risque maîtrisé par l'adoption graduelle et l'audit inaltérable."),
        ("03", "Sérénité",
         "Les dirigeants disposent d'une preuve chiffrée, auditable et exécutable de leur conformité."),
    ]
    cw = Inches(4.00); ch = Inches(1.90)
    cl = Inches(0.55); cg = Inches(0.30)
    for i, (num, h, d) in enumerate(pillars):
        x = cl + (cw + cg) * i
        y = Inches(4.60)
        add_panel(s, x, y, cw, ch)
        add_text(s, x + Inches(0.25), y + Inches(0.20),
                 Inches(0.8), Inches(0.50),
                 num, size=22, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x + Inches(1.0), y + Inches(0.25),
                 cw - Inches(1.3), Inches(0.50),
                 h, size=16, bold=True, color=CREAM)
        add_text(s, x + Inches(0.25), y + Inches(0.95),
                 cw - Inches(0.5), Inches(0.90),
                 d, size=11.5, color=CREAM_MUTE, line_spacing=1.30)
    # Next step bar
    add_panel(s, Inches(0.55), Inches(6.65),
              Inches(12.23), Inches(0.40),
              fill=INVIVOO_BLUE, border=INVIVOO_BLUE, border_weight=0)
    add_text(s, Inches(0.55), Inches(6.65),
             Inches(12.23), Inches(0.40),
             "PROCHAINE ÉTAPE  ·  Lancer un POC AppControl sur une première application — sous 6 semaines.",
             size=12, bold=True, color=CREAM,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, n)

# ── Build the deck ───────────────────────────────────────────────────
# Pair (illustration, content_builder).
plan = [
    (f"{ILLUSTRATIONS}/01-cover.png",       None),                    # 1 + 2 (cover x2)
    (f"{ILLUSTRATIONS}/02-probleme.png",    slide_probleme),          # 3 + 4
    (f"{ILLUSTRATIONS}/03-dora.png",        slide_dora),              # 5 + 6
    (f"{ILLUSTRATIONS}/04-cmdb.png",        slide_piege),             # 7 + 8
    (f"{ILLUSTRATIONS}/05-valeur.png",      slide_valeur),            # 9 + 10
    (f"{ILLUSTRATIONS}/06-xray.png",        slide_phase1),            # 11 + 12
    (f"{ILLUSTRATIONS}/07-telecommande.png", slide_phase2),           # 13 + 14
    (f"{ILLUSTRATIONS}/08-rebuild.png",     slide_phase3),            # 15 + 16
    (f"{ILLUSTRATIONS}/09-escalier.png",    slide_escalier),          # 17 + 18
    (f"{ILLUSTRATIONS}/10-gouvernance.png", slide_gouvernance),       # 19 + 20
    (f"{ILLUSTRATIONS}/11-reseau.png",      slide_reseau),            # 21 + 22
    (f"{ILLUSTRATIONS}/12-adoption.png",    slide_adoption),          # 23 + 24
    (f"{ILLUSTRATIONS}/13-conclusion.png",  slide_conclusion),        # 25 + 26
]

# Special handling for cover: the V1 illustration #1 IS the cover.
# Then we add a "dark navy" cover slide (the original V1 slide 2) re-skinned.
def slide_cover_dark(n):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    # Decorative left band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              0, 0, Inches(0.30), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = INVIVOO_BLUE
    band.line.fill.background()
    add_text(s, Inches(0.80), Inches(0.70), Inches(8), Inches(0.4),
             "RÉSILIENCE OPÉRATIONNELLE  ·  DORA",
             size=11, bold=True, color=INVIVOO_BLUE_2)
    add_text(s, Inches(0.80), Inches(1.30), Inches(12), Inches(2.4),
             ["Répondre avec sérénité",
              "aux exigences de DORA"],
             size=58, bold=True, color=CREAM, line_spacing=1.05)
    add_text(s, Inches(0.80), Inches(4.20), Inches(11), Inches(1.4),
             ["Industrialiser la résilience opérationnelle",
              "et prouver la conformité réglementaire —",
              "sans tout redocumenter, sans tout reconstruire."],
             size=20, color=CREAM_MUTE, line_spacing=1.35)
    indicators = [
        ("01", "NŒUD DE DONNÉES",  "Sécurisé"),
        ("02", "CONFORMITÉ",       "Activée"),
        ("03", "ARCHITECTURE",     "Résiliente"),
        ("04", "FLUX OPÉRATIONNEL","Unifié"),
    ]
    w = Inches(2.80); gap = Inches(0.20); left0 = Inches(0.80)
    for i, (n_, label, value) in enumerate(indicators):
        x = left0 + (w + gap) * i
        add_text(s, x, Inches(6.10), w, Inches(0.30),
                 n_, size=10, bold=True, color=INVIVOO_BLUE_2)
        add_text(s, x, Inches(6.40), w, Inches(0.30),
                 label, size=10, bold=True, color=CREAM_MUTE)
        add_text(s, x, Inches(6.70), w, Inches(0.30),
                 value, size=14, bold=True, color=CREAM)
        add_rule(s, x, Inches(7.00), Inches(2.40),
                 color=INVIVOO_BLUE, weight=1.5)
    add_text(s, Inches(0.80), Inches(7.20), Inches(8), Inches(0.30),
             "Invivoo  ·  Mai 2026  ·  Pitch AppControl",
             size=10, color=CREAM_MUTE)
    add_text(s, Inches(11.0), Inches(7.20), Inches(1.8), Inches(0.30),
             f"{n} / 26", size=10, color=CREAM_MUTE,
             align=PP_ALIGN.RIGHT)

plan[0] = (plan[0][0], slide_cover_dark)

n = 0
for image_path, content_builder in plan:
    n += 1
    illustration_slide(image_path)
    if content_builder is not None:
        n += 1
        content_builder(n)

OUT = "/tmp/decks/AppControl_DORA_Invivoo_Dark.pptx"
prs.save(OUT)
print(f"Saved {OUT} — {len(prs.slides)} slides")
